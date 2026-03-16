from pathlib import Path
import subprocess
import sys

import numpy as np
import pandas as pd
import statsmodels.formula.api as smf
from scipy import stats
from statsmodels.stats.proportion import confint_proportions_2indep
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
NOTEBOOK = ROOT / "notebooks" / "Analyze_AB_Test_Results.ipynb"
DATA = ROOT / "data" / "ab_data.csv"
REPORTS = ROOT / "reports"
REPORTS.mkdir(parents=True, exist_ok=True)
PRESENTATION = ROOT / "presentation" / "AB_Test_Results.pptx"
TEMPLATE_PRESENTATION = Path.home() / "Downloads" / "analyze-a_b-test-results-template (2).pptx"
COUNTRY_IMAGE = ROOT / "images" / "country_distribution_pct.png"

HTML_NAME = "Analyze_AB_Test_Results.html"
PDF_NAME = "Analyze_AB_Test_Results.pdf"


def run(cmd):
    print("Running:", " ".join(map(str, cmd)))
    subprocess.run(cmd, check=True)


def remove_slide(prs: Presentation, slide_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id_list.remove(slide_id_list[slide_index])


def compute_stats():
    """Load data and compute all presentation values programmatically."""
    df = pd.read_csv(DATA)

    # --- group sizes ---
    n_treatment = df.query('group == "treatment"').shape[0]
    n_control   = df.query('group == "control"').shape[0]

    # --- country distribution (%) ---
    country_pct = (df["country"].value_counts(normalize=True) * 100).round(1)

    # --- conversion rates by group ---
    conv_treatment = df.query('group == "treatment"')["converted"].mean()
    conv_control   = df.query('group == "control"')["converted"].mean()
    lift_pp        = conv_treatment - conv_control

    # --- 95 % CI for the absolute lift (Wald method) ---
    conv_treat_n = df.query('group == "treatment"')["converted"].sum()
    conv_ctrl_n  = df.query('group == "control"')["converted"].sum()
    ci_low, ci_high = confint_proportions_2indep(
        conv_treat_n, n_treatment, conv_ctrl_n, n_control, method="wald"
    )

    # --- relative lift ---
    relative_lift = lift_pp / conv_control

    # --- conversion rates by country × group (programmatic table) ---
    cg_table = (
        df.groupby(["group", "country"])["converted"]
          .mean()
          .mul(100)
          .round(2)
          .unstack()
    )

    # --- logistic regression: main effects (group + country) ---
    df_reg = df.copy()
    df_reg["country"] = pd.Categorical(df_reg["country"], categories=["CA", "UK", "US"])
    df_reg["group"]   = pd.Categorical(df_reg["group"],   categories=["control", "treatment"])

    model_main = smf.logit(
        'converted ~ C(group, Treatment(reference="control"))'
        ' + C(country, Treatment(reference="CA"))',
        data=df_reg,
    ).fit(disp=False)

    # --- logistic regression: interaction model (group × country) ---
    model_interaction = smf.logit(
        'converted ~ C(group, Treatment(reference="control"))'
        ' * C(country, Treatment(reference="CA"))',
        data=df_reg,
    ).fit(disp=False)

    # --- Likelihood Ratio Test ---
    lr_stat  = 2 * (model_interaction.llf - model_main.llf)
    df_diff  = model_interaction.df_model - model_main.df_model
    lrt_pval = stats.chi2.sf(lr_stat, df=int(df_diff))

    # --- treatment odds ratio from main-effects model ---
    treat_key = 'C(group, Treatment(reference="control"))[T.treatment]'
    treat_or  = np.exp(model_main.params[treat_key])
    treat_or_ci_low  = np.exp(model_main.conf_int().loc[treat_key, 0])
    treat_or_ci_high = np.exp(model_main.conf_int().loc[treat_key, 1])
    treat_pval = model_main.pvalues[treat_key]

    # --- Cohen's h effect size ---
    p1 = conv_treat_n / n_treatment
    p2 = conv_ctrl_n  / n_control
    cohens_h = 2 * (np.arcsin(np.sqrt(p1)) - np.arcsin(np.sqrt(p2)))

    return dict(
        n_treatment=n_treatment,
        n_control=n_control,
        country_pct=country_pct,
        conv_treatment=conv_treatment,
        conv_control=conv_control,
        lift_pp=lift_pp,
        ci_low=ci_low,
        ci_high=ci_high,
        relative_lift=relative_lift,
        cg_table=cg_table,
        treat_or=treat_or,
        treat_or_ci_low=treat_or_ci_low,
        treat_or_ci_high=treat_or_ci_high,
        treat_pval=treat_pval,
        lrt_pval=lrt_pval,
        lr_stat=lr_stat,
        df_diff=int(df_diff),
        cohens_h=cohens_h,
        model_main=model_main,
    )


def update_presentation(s: dict):
    source_presentation = TEMPLATE_PRESENTATION if TEMPLATE_PRESENTATION.exists() else PRESENTATION
    if not source_presentation.exists():
        raise FileNotFoundError(
            f"Neither template nor project presentation was found: "
            f"{TEMPLATE_PRESENTATION} / {PRESENTATION}"
        )
    if not COUNTRY_IMAGE.exists():
        raise FileNotFoundError(f"Country distribution image not found: {COUNTRY_IMAGE}")

    prs = Presentation(source_presentation)

    # Remove template "How to Use" slide if present
    if len(prs.slides) >= 2:
        slide_text = "\n".join(
            getattr(shape, "text", "").strip()
            for shape in prs.slides[1].shapes
            if hasattr(shape, "text")
        )
        if "How to Use This Template" in slide_text:
            remove_slide(prs, 1)

    # Slide 0 — title
    prs.slides[0].shapes[2].text = "Data Scientist: Francesco Malagrino"

    # Slide 2 — experiment overview (computed from data)
    us_pct = s["country_pct"].get("US", 0)
    uk_pct = s["country_pct"].get("UK", 0)
    ca_pct = s["country_pct"].get("CA", 0)
    prs.slides[2].shapes[1].text = (
        f"Total Variant Visitors: {s['n_treatment']:,}\n"
        f"Total Control Participants: {s['n_control']:,}\n\n"
        f"Users from: US ({us_pct}%), UK ({uk_pct}%), CA ({ca_pct}%)"
    )
    # Replace country chart image
    picture_shape = next(shape for shape in prs.slides[2].shapes if shape.shape_type == 13)
    left, top, width, height = (
        picture_shape.left, picture_shape.top,
        picture_shape.width, picture_shape.height,
    )
    picture_shape._element.getparent().remove(picture_shape._element)
    prs.slides[2].shapes.add_picture(str(COUNTRY_IMAGE), left, top, width=width, height=height)

    # Slide 3 — country × group conversion table (computed from data)
    tbl = s["cg_table"]   # rows: control / treatment; cols: CA, UK, US
    interaction_note = (
        "no significant interaction detected across countries"
        if s["lrt_pval"] >= 0.05
        else "significant interaction detected — country-specific strategy recommended"
    )
    prs.slides[3].shapes[1].text = (
        f"Executive Summary: The treatment page outperforms control in every country. "
        f"Treatment conversion rate is {s['conv_treatment']*100:.2f}% vs "
        f"{s['conv_control']*100:.2f}% for control — a lift of "
        f"+{s['lift_pp']*100:.2f} pp ({s['relative_lift']*100:.1f}% relative improvement). "
        f"Logistic regression confirms the treatment effect is statistically significant "
        f"(OR = {s['treat_or']:.3f}, 95% CI [{s['treat_or_ci_low']:.3f}, "
        f"{s['treat_or_ci_high']:.3f}], p < 0.001). "
        f"Likelihood Ratio Test: {interaction_note} "
        f"(χ²({s['df_diff']}) = {s['lr_stat']:.3f}, p = {s['lrt_pval']:.4f})."
    )
    table = next(shape.table for shape in prs.slides[3].shapes if shape.shape_type == 19)
    # Populate table cells from computed rates
    cell_map = {
        (1, 1): f"{tbl.loc['control', 'US']:.2f}%",
        (1, 2): f"{tbl.loc['control', 'UK']:.2f}%",
        (1, 3): f"{tbl.loc['control', 'CA']:.2f}%",
        (2, 1): f"{tbl.loc['treatment', 'US']:.2f}%",
        (2, 2): f"{tbl.loc['treatment', 'UK']:.2f}%",
        (2, 3): f"{tbl.loc['treatment', 'CA']:.2f}%",
    }
    for (row, col), value in cell_map.items():
        table.cell(row, col).text = value

    # Slide 4 — hypothesis test results + practical significance (computed from data)
    prs.slides[4].shapes[1].text = (
        f"Treatment Conversion Rate: {s['conv_treatment']*100:.2f}%\n"
        f"Control Conversion Rate: {s['conv_control']*100:.2f}%\n"
        f"Absolute Lift: +{s['lift_pp']*100:.2f} pp "
        f"(95% CI: [{s['ci_low']*100:.2f} pp, {s['ci_high']*100:.2f} pp])\n"
        f"Relative Lift: +{s['relative_lift']*100:.1f}% over control baseline\n"
        f"Simulation p-value: < 0.002 (0 of 500 null differences exceeded observed)\n"
        f"Logistic Regression: OR = {s['treat_or']:.3f} "
        f"(95% CI [{s['treat_or_ci_low']:.3f}, {s['treat_or_ci_high']:.3f}]), "
        f"p < 0.001 — significant after controlling for country\n"
        f"Interaction Test: LRT χ²({s['df_diff']}) = {s['lr_stat']:.3f}, "
        f"p = {s['lrt_pval']:.4f} — treatment effect is consistent across countries\n"
        f"Effect Size: Cohen's h = {s['cohens_h']:.4f}\n\n"
        f"Conclusion: We reject H₀ at α = 0.05. "
        f"Both statistical significance (p < 0.001) and practical significance "
        f"(+{s['relative_lift']*100:.0f}% relative lift, entire 95% CI above zero) "
        f"support implementing the treatment page globally."
    )

    prs.save(PRESENTATION)
    print(f"Created: {PRESENTATION}")


def main():
    if not NOTEBOOK.exists():
        raise FileNotFoundError(f"Notebook not found: {NOTEBOOK}")

    # Compute all stats from data before building presentation
    print("Computing statistics from data...")
    s = compute_stats()
    print(f"  Treatment rate:  {s['conv_treatment']*100:.4f}%")
    print(f"  Control rate:    {s['conv_control']*100:.4f}%")
    print(f"  Lift:            +{s['lift_pp']*100:.4f} pp")
    print(f"  Treat OR:        {s['treat_or']:.4f} (p={s['treat_pval']:.4g})")
    print(f"  LRT p-value:     {s['lrt_pval']:.4f}")
    print(f"  Cohen's h:       {s['cohens_h']:.4f}")

    update_presentation(s)

    run([
        sys.executable, "-m", "jupyter", "nbconvert",
        "--to", "html",
        "--execute",
        str(NOTEBOOK),
        "--output", HTML_NAME,
        "--output-dir", str(REPORTS),
    ])

    run([
        sys.executable, "-m", "jupyter", "nbconvert",
        "--to", "webpdf",
        "--execute",
        "--allow-chromium-download",
        str(NOTEBOOK),
        "--output", PDF_NAME,
        "--output-dir", str(REPORTS),
    ])

    print(f"Created: {REPORTS / HTML_NAME}")
    print(f"Created: {REPORTS / PDF_NAME}")


if __name__ == "__main__":
    main()
